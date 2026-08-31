import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# Color Scheme (Cybersecurity Dark Theme)
BG_COLOR = RGBColor(17, 22, 34)       # #111622 Deep Dark Slate
TEXT_WHITE = RGBColor(248, 250, 252)  # #f8fafc Off-White
TEXT_MUTED = RGBColor(148, 163, 184)  # #94a3b8 Muted Grey
ACCENT_BLUE = RGBColor(56, 189, 248)  # #38bdf8 Neon Blue
ACCENT_GREEN = RGBColor(16, 185, 129) # #10b981 Emerald Green
ACCENT_RED = RGBColor(239, 68, 68)    # #ef4444 Danger Red

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="SENTINELGATE // THREAT MANAGEMENT"):
    # Add top category label
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Consolas'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE

    # Add main slide title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

# ==============================================================================
# SLIDE 1: Title Slide (Dark Cyber Splash)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank Layout
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

# Large Title text box
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True

# Main title
p1 = tf.paragraphs[0]
p1.text = "🛡️ SENTINELGATE"
p1.font.name = 'Segoe UI'
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = ACCENT_BLUE
p1.space_after = Pt(10)

# Subtitle
p2 = tf.add_paragraph()
p2.text = "Adaptive Rate-Limiting & Abuse Detection API Gateway"
p2.font.name = 'Segoe UI'
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_WHITE
p2.space_after = Pt(20)

# Extra description
p3 = tf.add_paragraph()
p3.text = "Next-generation API shield combining distributed Redis Token Buckets with live Machine Learning diagnostics."
p3.font.name = 'Segoe UI'
p3.font.size = Pt(14)
p3.font.italic = True
p3.font.color.rgb = TEXT_MUTED

# ==============================================================================
# SLIDE 2: The Security Problem (Static Limits vs Scrapers)
# ==============================================================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_header(slide2, "The Vulnerability of Static Rate Limits", "PROBLEM SPECIFICATION")

# Left Column (Core issue description)
left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
tf_l = left_box.text_frame
tf_l.word_wrap = True

p_l1 = tf_l.paragraphs[0]
p_l1.text = "Why Traditional Gateways Fail"
p_l1.font.name = 'Segoe UI'
p_l1.font.size = Pt(18)
p_l1.font.bold = True
p_l1.font.color.rgb = ACCENT_RED
p_l1.space_after = Pt(15)

p_l2 = tf_l.add_paragraph()
p_l2.text = "Static rate-limiting tools apply uniform thresholds (e.g. 60 requests/minute) globally to all clients. This architecture creates a critical security mismatch:\n\n" \
            "• Honest, bursty users are locked out during normal burst load spikes, creating high false-positive rates.\n" \
            "• Malicious actors easily bypass limits using low-and-slow scraping, bot rotations, or DDoS patterns.\n" \
            "• Real-time mitigation usually requires permanent manual IP blocks, leading to massive administrative overhead."
p_l2.font.name = 'Segoe UI'
p_l2.font.size = Pt(13)
p_l2.font.color.rgb = TEXT_MUTED

# Right Column (Detailed Impact Blocks)
right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
tf_r = right_box.text_frame
tf_r.word_wrap = True

def add_threat_block(tf, heading, text):
    h = tf.add_paragraph()
    h.text = f"🔥 {heading}"
    h.font.name = 'Segoe UI'
    h.font.size = Pt(15)
    h.font.bold = True
    h.font.color.rgb = TEXT_WHITE
    h.space_after = Pt(4)
    
    t = tf.add_paragraph()
    t.text = text
    t.font.name = 'Segoe UI'
    t.font.size = Pt(12)
    t.font.color.rgb = TEXT_MUTED
    t.space_after = Pt(15)

tf_r.paragraphs[0].text = "Sophisticated Attack Vectors Addressed:"
tf_r.paragraphs[0].font.name = 'Segoe UI'
tf_r.paragraphs[0].font.size = Pt(16)
tf_r.paragraphs[0].font.bold = True
tf_r.paragraphs[0].font.color.rgb = TEXT_WHITE
tf_r.paragraphs[0].space_after = Pt(12)

add_threat_block(tf_r, "Aggressive API Scraping", "Repetitive, high-volume endpoint traversal aiming to steal proprietary catalog or product data.")
add_threat_block(tf_r, "Credential Brute-Forcing", "Automated authentication attempts returning consistent 401 Unauthorized errors.")
add_threat_block(tf_r, "Distributed DDoS Bursts", "Sudden coordinate traffic spikes intended to degrade microservice availability.")

# ==============================================================================
# SLIDE 3: The Architecture (Token Bucket & Redis)
# ==============================================================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_header(slide3, "Distributed Token Bucket Architecture", "INFRASTRUCTURE LAYOUT")

# Left Column (System components)
left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
tf_l = left_box.text_frame
tf_l.word_wrap = True

p_l1 = tf_l.paragraphs[0]
p_l1.text = "FastAPI & Redis Distributed State"
p_l1.font.name = 'Segoe UI'
p_l1.font.size = Pt(18)
p_l1.font.bold = True
p_l1.font.color.rgb = ACCENT_BLUE
p_l1.space_after = Pt(15)

p_l2 = tf_l.add_paragraph()
p_l2.text = "Rather than relying on local memory limits, SentinelGate synchronizes state across multiple gateway nodes globally:\n\n" \
            "• Gateway Middleware: Lightweight FastAPI proxy nodes verifying requests in sub-milliseconds.\n" \
            "• Distributed Redis Store: Real-time atomic token evaluation using embedded Lua script execution.\n" \
            "• DB Logging Pipeline: Asynchronous SQLAlchemy task execution logging metrics without blocking core API requests."
p_l2.font.name = 'Segoe UI'
p_l2.font.size = Pt(13)
p_l2.font.color.rgb = TEXT_MUTED

# Right Column (Technical concepts explained)
right_box = slide3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
tf_r = right_box.text_frame
tf_r.word_wrap = True

tf_r.paragraphs[0].text = "Distributed Engineering Terms:"
tf_r.paragraphs[0].font.name = 'Segoe UI'
tf_r.paragraphs[0].font.size = Pt(16)
tf_r.paragraphs[0].font.bold = True
tf_r.paragraphs[0].font.color.rgb = TEXT_WHITE
tf_r.paragraphs[0].space_after = Pt(12)

add_threat_block(tf_r, "Token Bucket Algorithm", "An atomic algorithm tracking available 'tokens' which refill at a set rate. Capacity and refill rate scale dynamically with risk.")
add_threat_block(tf_r, "Lua Scripting Atomicity", "Executing token calculations inside Redis in a single atomic script to prevent race conditions across parallel requests.")
add_threat_block(tf_r, "Resilient Fallbacks", "Integrated auto-fallback logic immediately switching to InMemoryRedis and SQLite if production infrastructure drops.")

# ==============================================================================
# SLIDE 4: The Intelligence Engine (ML & Scoring)
# ==============================================================================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_header(slide4, "Heuristic Risk Analysis & Isolation Forest ML", "INTELLIGENCE ENGINE")

# Left Column (Telemetry features)
left_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
tf_l = left_box.text_frame
tf_l.word_wrap = True

p_l1 = tf_l.paragraphs[0]
p_l1.text = "Multi-Factor Threat Assessment"
p_l1.font.name = 'Segoe UI'
p_l1.font.size = Pt(18)
p_l1.font.bold = True
p_l1.font.color.rgb = ACCENT_GREEN
p_l1.space_after = Pt(15)

p_l2 = tf_l.add_paragraph()
p_l2.text = "Clients are continuously monitored in sliding time windows to compute behavior diagnostics:\n\n" \
            "• Request Rate (RPM): Current volume compared to safe baseline thresholds.\n" \
            "• Burstiness (Variance): Time delta dispersion tracking machine-like rapid calls.\n" \
            "• Error Rate Frequency: Percentage of unauthorized or broken endpoints reached.\n" \
            "• Endpoint Repetition: Ratio of repeating targets hit, flagging scraping bots."
p_l2.font.name = 'Segoe UI'
p_l2.font.size = Pt(13)
p_l2.font.color.rgb = TEXT_MUTED

# Right Column (Isolation Forest & Scoring)
right_box = slide4.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
tf_r = right_box.text_frame
tf_r.word_wrap = True

tf_r.paragraphs[0].text = "Machine Learning & Adaptive Control:"
tf_r.paragraphs[0].font.name = 'Segoe UI'
tf_r.paragraphs[0].font.size = Pt(16)
tf_r.paragraphs[0].font.bold = True
tf_r.paragraphs[0].font.color.rgb = TEXT_WHITE
tf_r.paragraphs[0].space_after = Pt(12)

add_threat_block(tf_r, "Isolation Forest Classifier", "Unsupervised machine learning model that isolates anomalous vectors without pre-labeled data. Highly sensitive to novel bot patterns.")
add_threat_block(tf_r, "Explainable Risk Mapping (0-100)", "Heuristic flags combined with ML anomalies yield a score translating to: NORMAL (100 req/m) ➡️ MONITORED (60 req/m) ➡️ THROTTLED (10 req/m) ➡️ BLOCKED (0 req/m).")

# ==============================================================================
# SLIDE 5: Business Value & Key Takeaways
# ==============================================================================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_header(slide5, "Operational Advantages & Demo Value", "KEY TAKEAWAYS")

# Left Column (Business & Technical Value)
left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
tf_l = left_box.text_frame
tf_l.word_wrap = True

p_l1 = tf_l.paragraphs[0]
p_l1.text = "Why SentinelGate Wins"
p_l1.font.name = 'Segoe UI'
p_l1.font.size = Pt(18)
p_l1.font.bold = True
p_l1.font.color.rgb = ACCENT_BLUE
p_l1.space_after = Pt(15)

p_l2 = tf_l.add_paragraph()
p_l2.text = "• Automatic Recovery: Risk scores decay dynamically when abuse stops, preventing permanent lockouts.\n" \
            "• Minimal Overhead: In-memory evaluation path maintains a sub-millisecond check latency footprint.\n" \
            "• Real-Time Control Link: WebSocket telemetry pushes logs and metrics instantly to an operational dashboard.\n" \
            "• Fully Deployable Stack: Built-in sandbox simulator generated out-of-the-box (Normal, DDoS, Brute Force attacks)."
p_l2.font.name = 'Segoe UI'
p_l2.font.size = Pt(13)
p_l2.font.color.rgb = TEXT_MUTED

# Right Column (Conclusion / Demo Guide)
right_box = slide5.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
tf_r = right_box.text_frame
tf_r.word_wrap = True

tf_r.paragraphs[0].text = "The Presentation Storyline:"
tf_r.paragraphs[0].font.name = 'Segoe UI'
tf_r.paragraphs[0].font.size = Pt(16)
tf_r.paragraphs[0].font.bold = True
tf_r.paragraphs[0].font.color.rgb = TEXT_WHITE
tf_r.paragraphs[0].space_after = Pt(12)

add_threat_block(tf_r, "Step 1: Establishing the Baseline", "Show legitimate normal traffic passing seamlessly on the dashboard telemetry chart.")
add_threat_block(tf_r, "Step 2: Triggering the Attack & Defending", "Start a DDoS burst. Observe the threat indexes scale dynamically to automatically throttle and isolate the malicious client.")
add_threat_block(tf_r, "Step 3: Self-Healing", "Halt the attack. Observe the risk decaying automatically and restoring client access safely.")

# Save presentation
prs.save("SentinelGate_Presentation.pptx")
print("PowerPoint presentation generated successfully!")
